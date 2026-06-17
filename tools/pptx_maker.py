"""python-pptx 기반 PPT 생성기.
디자인 기준: skills/pptx/design_spec.md (흑백, Pretendard 폰트, 미니멀)
"""

import math
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ── 색상 팔레트 (흑백 전용) ──────────────────────────────────────
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY_DARK = RGBColor(0x22, 0x22, 0x22)
GRAY_MID = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
GRAY_XLIGHT = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ── 폰트 우선순위 (Pretendard 없으면 시스템 폰트 폴백) ───────────
FONT_PRIMARY = "Pretendard"

# ── 레이아웃 상수 (16:9 와이드스크린 표준, 인치) ────────────────
SLIDE_W = 13.33
SLIDE_H = 7.5
MARGIN = 0.6
CONTENT_W = SLIDE_W - MARGIN * 2
BOX_PAD = 0.18  # 박스 안 텍스트 프레임 내부 여백 (글자가 테두리에 붙지 않게)

OUTPUT_DIR = Path(__file__).parent.parent / "data" / "outputs"


def _set_font(run, size_pt: int, bold: bool = False, color: RGBColor = GRAY_DARK) -> None:
    run.font.name = FONT_PRIMARY
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def _pad_tf(tf, pad: float = BOX_PAD, top: float | None = None) -> None:
    """텍스트 프레임 내부 여백을 준다 — 글자가 박스 테두리에 붙는 답답함 해소."""
    tf.margin_left = Inches(pad)
    tf.margin_right = Inches(pad)
    tf.margin_top = Inches(top if top is not None else pad * 0.6)
    tf.margin_bottom = Inches(top if top is not None else pad * 0.6)


def _est_lines(text: str, box_w_in: float, size_pt: int,
               pad: float = BOX_PAD) -> int:
    """주어진 너비/폰트에서 텍스트가 몇 줄로 줄바꿈될지 추정한다.

    한글 한 글자 폭 ≈ size_pt/72 인치로 보고(영문/공백은 더 좁지만 보수적으로
    잡아 넘침을 막는다), 사용 가능한 폭으로 나눈다.
    """
    usable = max(0.5, box_w_in - 2 * pad)
    char_w = size_pt * 0.95 / 72.0
    per_line = max(1, int(usable / char_w))
    return max(1, math.ceil(len(text) / per_line))


def _fit_font_size(text: str, box_w_in: float, box_h_in: float,
                   max_size: int, min_size: int = 14,
                   pad: float = BOX_PAD) -> int:
    """텍스트가 박스(w×h)를 넘지 않는 최대 폰트 크기를 찾는다.

    줄바꿈으로 늘어나는 줄 수를 고려해 크기를 1pt씩 낮춰가며 맞춘다.
    python-pptx는 실측 폭을 모르므로 보수적 휴리스틱(넘침보다 약간 작게).
    """
    text = text or ""
    usable_h = max(0.2, box_h_in - 2 * (pad * 0.6))
    for size in range(max_size, min_size - 1, -1):
        lines = _est_lines(text, box_w_in, size, pad)
        need_h = lines * size * 1.25 / 72.0
        if need_h <= usable_h:
            return size
    return min_size


def _add_divider_line(slide, y: float) -> None:
    shape = slide.shapes.add_shape(
        1,
        Inches(MARGIN), Inches(y),
        Inches(CONTENT_W), Inches(0),
    )
    shape.line.color.rgb = GRAY_LIGHT
    shape.line.width = Pt(1)


def _add_text_box(slide, text: str, x, y, w, h,
                  size: int, bold: bool = False,
                  color: RGBColor = GRAY_DARK,
                  align=PP_ALIGN.LEFT) -> None:
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    # 명시적 줄바꿈(\n)은 별도 단락으로 렌더 — 단일 run에 넣으면 무시되므로.
    for idx, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        _set_font(run, size, bold, color)


def _cell_font(cell, size_pt: int, bold: bool = False,
               color: RGBColor = GRAY_DARK, align=PP_ALIGN.LEFT) -> None:
    tf = cell.text_frame
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            _set_font(run, size_pt, bold, color)
        if not para.runs:
            run = para.add_run()
            _set_font(run, size_pt, bold, color)


def _distributed_bullets(slide, items, x, y, w, h,
                         base_size: int, marker: str = "▪  ",
                         color: RGBColor = GRAY_DARK,
                         numbered: bool = False) -> None:
    """불릿/번호 목록을 주어진 영역(x,y,w,h) 세로로 골고루 분포시킨다.

    - 항목이 적으면 폰트를 키우고 간격을 넓혀 빈 여백을 채운다.
    - 텍스트 프레임을 영역 전체로 잡고 항목 사이 간격을 동적 계산한다.
    """
    n = len(items)
    if n == 0:
        return

    # 항목 수에 따라 폰트 동적 조정 (적을수록 크게)
    if n <= 2:
        size = base_size + 5
    elif n == 3:
        size = base_size + 3
    elif n == 4:
        size = base_size + 1
    else:
        size = base_size

    # 가장 긴 항목 기준으로 폰트를 줄여 텍스트가 영역을 넘지 않게 자동 맞춤.
    marker_len = 3 if not numbered else 4
    longest = max((len(it) for it in items), default=0) + marker_len
    avg_lines = sum(_est_lines(("X" * marker_len) + it, w, size) for it in items)
    line_unit = size * 1.25 / 72
    while size > base_size - 6 and avg_lines * line_unit > h * 0.92:
        size -= 1
        line_unit = size * 1.25 / 72
        avg_lines = sum(_est_lines(("X" * marker_len) + it, w, size) for it in items)

    # 줄바꿈으로 늘어난 실제 줄 수를 반영해 빈 공간을 항목 사이에 분배.
    text_total = avg_lines * line_unit
    free = max(0.0, h - text_total)
    gap_in = free / (n + 1)
    gap_pt = max(8.0, min(48.0, gap_in * 72))

    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    _pad_tf(tf, pad=0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # 세로 중앙 기준 + 간격으로 분포

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = (f"{i + 1}.  " if numbered else marker) + item
        _set_font(run, size, color=color)
        p.space_before = Pt(gap_pt)
        p.space_after = Pt(0)
        p.line_spacing = 1.2


class PptxMaker:
    """흑백 미니멀 스타일 PPT 생성기 (16:9 와이드스크린)."""

    def __init__(self, title: str = "프레젠테이션") -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W)
        self.prs.slide_height = Inches(SLIDE_H)
        self.prs.core_properties.title = title
        self.prs.core_properties.author = "AI 교육팀"
        self._blank_layout = self.prs.slide_layouts[6]
        self.footer_label = title
        self._page = 0

    def _new_slide(self):
        slide = self.prs.slides.add_slide(self._blank_layout)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = WHITE
        return slide

    def _add_footer(self, slide) -> None:
        """콘텐츠 슬라이드 우하단에 '커리큘럼명 · 페이지' 러닝 푸터를 단다.

        흑백 톤을 유지하며 슬라이드 간 통일감을 준다(표지/구분 슬라이드 제외).
        """
        self._page += 1
        _add_text_box(slide, f"{self.footer_label}  ·  {self._page:02d}",
                      x=MARGIN, y=SLIDE_H - 0.5, w=CONTENT_W, h=0.35,
                      size=11, color=GRAY_MID, align=PP_ALIGN.RIGHT)

    def add_title_slide(self, title: str, subtitle: str = "") -> None:
        """타이틀 슬라이드 — 중앙 정렬, 검정 배경."""
        slide = self._new_slide()

        # 검정 배경
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = BLACK

        # 제목
        _add_text_box(slide, title,
                      x=1.5, y=2.2, w=10.33, h=1.6,
                      size=44, bold=True, color=WHITE,
                      align=PP_ALIGN.CENTER)

        # 구분선
        shape = slide.shapes.add_shape(
            1,
            Inches(5.0), Inches(4.1),
            Inches(3.33), Inches(0),
        )
        shape.line.color.rgb = GRAY_MID
        shape.line.width = Pt(1)

        # 부제목
        if subtitle:
            _add_text_box(slide, subtitle,
                          x=1.5, y=4.3, w=10.33, h=1.4,
                          size=20, color=GRAY_MID,
                          align=PP_ALIGN.CENTER)

        # 우하단 라벨
        _add_text_box(slide, "AI 교육팀",
                      x=SLIDE_W - 2.0, y=SLIDE_H - 0.65, w=1.4, h=0.4,
                      size=13, color=GRAY_MID,
                      align=PP_ALIGN.RIGHT)

    def add_content_slide(self, section_label: str, title: str,
                          bullets: list[str]) -> None:
        """콘텐츠 슬라이드 — 섹션 라벨 + 제목 + 불릿."""
        slide = self._new_slide()

        # 상단 섹션 라벨 + 구분선
        _add_text_box(slide, section_label.upper(),
                      x=MARGIN, y=0.35, w=CONTENT_W, h=0.45,
                      size=13, color=GRAY_MID)
        _add_divider_line(slide, y=0.9)

        # 슬라이드 제목
        _add_text_box(slide, title,
                      x=MARGIN, y=1.0, w=CONTENT_W, h=1.0,
                      size=42, bold=True, color=BLACK)

        # 불릿 — 본문 영역 전체에 세로로 골고루 분포 (최대 5개)
        area_top = 2.25
        area_h = SLIDE_H - area_top - 0.7  # 하단 푸터 공간 확보
        _distributed_bullets(slide, bullets[:5],
                             x=MARGIN, y=area_top, w=CONTENT_W, h=area_h,
                             base_size=27)

        self._add_footer(slide)

    def add_concept_slide(self, section_label: str, term: str,
                          explain: str, analogy: str = "") -> None:
        """개념 설명 슬라이드 — 용어(제목) + 풀어 쓴 설명 + 일상 비유 박스.

        교재 원문 대신 비개발 실무자가 이해하기 쉬운 표현으로 설명한다.
        """
        slide = self._new_slide()
        title_size = 42 if len(term) <= 16 else (34 if len(term) <= 24 else 28)
        self._slide_header(slide, section_label, term, title_size=title_size)

        has_analogy = bool(analogy)
        area_top = 2.35
        # 비유 박스 자리를 아래에 확보하고 나머지를 설명 영역으로.
        analogy_h = 1.5 if has_analogy else 0.0
        analogy_y = SLIDE_H - 0.7 - analogy_h
        explain_h = (analogy_y - 0.25) - area_top

        ex_size = _fit_font_size(explain, CONTENT_W, explain_h,
                                 max_size=25, min_size=16, pad=0.05)
        eb = slide.shapes.add_textbox(Inches(MARGIN), Inches(area_top),
                                      Inches(CONTENT_W), Inches(explain_h))
        ef = eb.text_frame
        ef.word_wrap = True
        _pad_tf(ef, pad=0.05)
        ef.vertical_anchor = MSO_ANCHOR.TOP
        ep = ef.paragraphs[0]
        ep.alignment = PP_ALIGN.LEFT
        ep.line_spacing = 1.35
        er = ep.add_run()
        er.text = explain
        _set_font(er, ex_size, color=GRAY_DARK)

        if has_analogy:
            box = slide.shapes.add_shape(1, Inches(MARGIN), Inches(analogy_y),
                                         Inches(CONTENT_W), Inches(analogy_h))
            box.fill.solid()
            box.fill.fore_color.rgb = GRAY_XLIGHT
            box.line.color.rgb = GRAY_LIGHT
            box.line.width = Pt(1)
            box.shadow.inherit = False
            an_size = _fit_font_size("💡 비유   " + analogy, CONTENT_W, analogy_h,
                                     max_size=21, min_size=14, pad=0.2)
            tf = box.text_frame
            tf.word_wrap = True
            _pad_tf(tf, pad=0.25)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.25
            r1 = p.add_run()
            r1.text = "💡 비유   "
            _set_font(r1, an_size, bold=True, color=BLACK)
            r2 = p.add_run()
            r2.text = analogy
            _set_font(r2, an_size, color=GRAY_DARK)

        self._add_footer(slide)

    def _slide_header(self, slide, section_label: str, title: str,
                      title_size: int = 42) -> None:
        """섹션 라벨 + 구분선 + 큰 제목 (content/cards/flow 공통 상단)."""
        _add_text_box(slide, section_label.upper(),
                      x=MARGIN, y=0.35, w=CONTENT_W, h=0.45,
                      size=14, color=GRAY_MID)
        _add_divider_line(slide, y=0.9)
        _add_text_box(slide, title,
                      x=MARGIN, y=1.0, w=CONTENT_W, h=1.0,
                      size=title_size, bold=True, color=BLACK)

    def add_cards_slide(self, section_label: str, title: str,
                        items: list[str], variant: str = "number") -> None:
        """번호 카드 슬라이드 — 항목을 번호 배지 + 박스로 세로로 꽉 채운다.

        불릿보다 구조가 또렷하고 빈 여백을 채워 가독성이 좋다.
        """
        slide = self._new_slide()
        self._slide_header(slide, section_label, title)

        items = items[:5]
        n = len(items)
        if n == 0:
            return

        area_top = 2.35
        area_h = SLIDE_H - area_top - 0.7  # 하단 푸터 공간 확보
        gap = 0.22
        box_h = (area_h - gap * (n - 1)) / n
        base_txt = 28 if n <= 2 else (25 if n == 3 else (22 if n == 4 else 19))
        badge_w = 1.0
        text_w = CONTENT_W - badge_w - 0.6
        # 모든 카드가 잘리지 않도록 가장 긴 항목 기준 폰트를 통일해서 자동 축소.
        txt_size = min(
            _fit_font_size(it, text_w, box_h, max_size=base_txt, min_size=15)
            for it in items
        )
        num_size = txt_size + 10

        for i, item in enumerate(items):
            y = area_top + i * (box_h + gap)
            # 배경 박스 (연회색)
            box = slide.shapes.add_shape(1, Inches(MARGIN), Inches(y),
                                         Inches(CONTENT_W), Inches(box_h))
            box.fill.solid()
            box.fill.fore_color.rgb = GRAY_XLIGHT
            box.line.color.rgb = GRAY_LIGHT
            box.line.width = Pt(1)
            box.shadow.inherit = False
            # 번호 배지 (검정 정사각 왼쪽)
            badge = slide.shapes.add_shape(1, Inches(MARGIN), Inches(y),
                                           Inches(badge_w), Inches(box_h))
            badge.fill.solid()
            badge.fill.fore_color.rgb = BLACK
            badge.line.fill.background()
            badge.shadow.inherit = False
            bf = badge.text_frame
            bf.word_wrap = True
            bf.vertical_anchor = MSO_ANCHOR.MIDDLE
            bp = bf.paragraphs[0]
            bp.alignment = PP_ALIGN.CENTER
            br = bp.add_run()
            br.text = str(i + 1) if variant == "number" else "•"
            _set_font(br, num_size, bold=True, color=WHITE)
            # 항목 텍스트 (배지 오른쪽, 세로중앙)
            tb = slide.shapes.add_textbox(
                Inches(MARGIN + badge_w + 0.3), Inches(y),
                Inches(text_w), Inches(box_h))
            tf = tb.text_frame
            tf.word_wrap = True
            _pad_tf(tf, pad=0.14)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tp = tf.paragraphs[0]
            tp.alignment = PP_ALIGN.LEFT
            tp.line_spacing = 1.15
            tr = tp.add_run()
            tr.text = item
            _set_font(tr, txt_size, color=GRAY_DARK)

        self._add_footer(slide)

    def add_flow_slide(self, section_label: str, title: str,
                       steps: list[dict]) -> None:
        """구조도 슬라이드 — 3박스 가로 흐름(목표 ▶ 핵심 개념 ▶ 실습)."""
        slide = self._new_slide()
        self._slide_header(slide, section_label, title)

        steps = steps[:3]
        n = len(steps)
        if n == 0:
            return

        area_top = 2.45
        area_h = SLIDE_H - area_top - 0.6
        arrow_w = 0.7
        box_w = (CONTENT_W - arrow_w * (n - 1)) / n

        for i, step in enumerate(steps):
            x = MARGIN + i * (box_w + arrow_w)
            box = slide.shapes.add_shape(1, Inches(x), Inches(area_top),
                                         Inches(box_w), Inches(area_h))
            box.fill.solid()
            box.fill.fore_color.rgb = GRAY_XLIGHT
            box.line.color.rgb = BLACK
            box.line.width = Pt(1.5)
            box.shadow.inherit = False
            # 라벨 (상단 굵게)
            _add_text_box(slide, step.get("label", ""),
                          x=x, y=area_top + 0.18, w=box_w, h=0.6,
                          size=23, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
            # 항목들 (남은 영역에 분포)
            _distributed_bullets(slide, step.get("items", [])[:3],
                                 x=x + 0.25, y=area_top + 1.0,
                                 w=box_w - 0.5, h=area_h - 1.2,
                                 base_size=16, marker="·  ", color=GRAY_DARK)
            # 화살표
            if i < n - 1:
                _add_text_box(slide, "▶",
                              x=x + box_w, y=area_top + area_h / 2 - 0.35,
                              w=arrow_w, h=0.7, size=30, bold=True,
                              color=GRAY_MID, align=PP_ALIGN.CENTER)

        self._add_footer(slide)

    def add_table_slide(self, section_label: str, title: str,
                        headers: list[str], rows: list[list[str]]) -> None:
        """테이블 슬라이드 — 헤더(검정) + 데이터 행(흰색/연회색 교차)."""
        slide = self._new_slide()

        _add_text_box(slide, section_label.upper(),
                      x=MARGIN, y=0.35, w=CONTENT_W, h=0.45,
                      size=13, color=GRAY_MID)
        _add_divider_line(slide, y=0.9)

        _add_text_box(slide, title,
                      x=MARGIN, y=1.0, w=CONTENT_W, h=0.8,
                      size=38, bold=True, color=BLACK)

        if not headers or not rows:
            return

        n_cols = min(len(headers), 5)
        n_rows = min(len(rows), 8) + 1  # +1 for header row
        table_top = 2.1
        # 본문 영역 세로를 꽉 채운다 (행이 적어도 빈 여백 없이 분포). 하단 푸터 공간 확보.
        table_area_h = SLIDE_H - table_top - 0.7

        table = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(MARGIN), Inches(table_top),
            Inches(CONTENT_W), Inches(table_area_h),
        ).table

        # 열 너비 균등 분배
        col_w_in = CONTENT_W / n_cols
        col_w = int(Inches(col_w_in))
        for j in range(n_cols):
            table.columns[j].width = col_w

        # 행 높이 분배 — 헤더는 고정, 데이터 행은 남는 높이를 균등 분배.
        # PowerPoint는 이 값을 최소 높이로 보고 내용이 길면 행을 늘리므로 잘리지 않는다.
        header_h = 0.6
        data_h = max(0.5, (table_area_h - header_h) / (n_rows - 1))
        table.rows[0].height = Inches(header_h)
        for i in range(1, n_rows):
            table.rows[i].height = Inches(data_h)

        # 가장 긴 셀이 열 폭 안에서 잘리지 않도록 본문 폰트를 자동 산출해 통일.
        n_data = n_rows - 1
        base_body = 24 if n_data <= 4 else (21 if n_data <= 6 else 18)
        longest = ""
        for row in rows[:n_data]:
            for c in row[:n_cols]:
                if len(str(c)) > len(longest):
                    longest = str(c)
        body_size = _fit_font_size(longest, col_w_in, data_h,
                                   max_size=base_body, min_size=13, pad=0.12)
        head_size = body_size + 2

        cell_margin = Inches(0.12)
        cell_margin_v = Inches(0.06)

        # 헤더 행 (검정 배경, 흰 텍스트)
        for j, hdr in enumerate(headers[:n_cols]):
            cell = table.cell(0, j)
            cell.text = hdr
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLACK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = cell_margin
            cell.margin_top = cell.margin_bottom = cell_margin_v
            _cell_font(cell, head_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # 데이터 행 (흰색 / 연회색 교차)
        for i, row in enumerate(rows[:n_rows - 1]):
            bg = WHITE if i % 2 == 0 else GRAY_XLIGHT
            for j in range(n_cols):
                cell = table.cell(i + 1, j)
                cell_text = row[j] if j < len(row) else ""
                cell.text = cell_text
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_left = cell.margin_right = cell_margin
                cell.margin_top = cell.margin_bottom = cell_margin_v
                _cell_font(cell, body_size, color=GRAY_DARK)

        self._add_footer(slide)

    def add_comparison_slide(self, title: str,
                             left_label: str, left_items: list[str],
                             right_label: str, right_items: list[str]) -> None:
        """비교 슬라이드 — 2열 구성."""
        slide = self._new_slide()

        _add_text_box(slide, title,
                      x=MARGIN, y=0.35, w=CONTENT_W, h=0.8,
                      size=38, bold=True, color=BLACK)
        _add_divider_line(slide, y=1.3)

        col_w = (CONTENT_W - 0.3) / 2
        left_x = MARGIN
        right_x = MARGIN + col_w + 0.3

        col_top = 1.6
        items_top = 2.2
        items_h = SLIDE_H - items_top - 0.55

        _add_text_box(slide, left_label,
                      x=left_x, y=col_top, w=col_w, h=0.55,
                      size=24, bold=True, color=BLACK)
        _distributed_bullets(slide, left_items[:6],
                             x=left_x, y=items_top, w=col_w, h=items_h,
                             base_size=22)

        sep = slide.shapes.add_shape(
            1,
            Inches(MARGIN + col_w + 0.1), Inches(col_top),
            Inches(0), Inches(items_h + (items_top - col_top)),
        )
        sep.line.color.rgb = GRAY_LIGHT
        sep.line.width = Pt(1)

        _add_text_box(slide, right_label,
                      x=right_x, y=col_top, w=col_w, h=0.55,
                      size=24, bold=True, color=BLACK)
        _distributed_bullets(slide, right_items[:6],
                             x=right_x, y=items_top, w=col_w, h=items_h,
                             base_size=22)

    def add_summary_slide(self, lessons: list[str], source: str = "") -> None:
        """요약 슬라이드 — 핵심 교훈 + 출처."""
        slide = self._new_slide()

        # 연한 회색 배경
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = GRAY_XLIGHT

        _add_text_box(slide, "커리큘럼 요약",
                      x=MARGIN, y=0.5, w=CONTENT_W, h=0.8,
                      size=40, bold=True, color=BLACK)
        _add_divider_line(slide, y=1.45)

        area_top = 1.65
        area_h = SLIDE_H - area_top - 0.75  # 하단 출처 공간 확보
        _distributed_bullets(slide, lessons[:8],
                             x=MARGIN, y=area_top, w=CONTENT_W, h=area_h,
                             base_size=25, numbered=True)

        if source:
            _add_text_box(slide, f"출처: {source}",
                          x=MARGIN, y=SLIDE_H - 0.65, w=CONTENT_W, h=0.4,
                          size=13, color=GRAY_MID)

    def add_part_divider_slide(self, label: str, title: str, weeks: str = "") -> None:
        """파트 구분 슬라이드 — 검정 배경 전면, 큰 'PART A' + 파트 제목.

        커리큘럼이 Part A/B로 나뉠 때 각 파트 시작에 삽입한다.
        주차 구분 슬라이드(흰 배경)와 시각적으로 대비시켜 큰 단락을 표시한다.
        """
        slide = self._new_slide()

        # 검정 배경
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = BLACK

        # 큰 파트 라벨 (예: PART A)
        _add_text_box(slide, label or "PART",
                      x=1.5, y=2.0, w=10.33, h=1.4,
                      size=76, bold=True, color=WHITE,
                      align=PP_ALIGN.CENTER)

        # 구분선
        shape = slide.shapes.add_shape(
            1,
            Inches(5.0), Inches(3.7),
            Inches(3.33), Inches(0),
        )
        shape.line.color.rgb = GRAY_MID
        shape.line.width = Pt(1)

        # 파트 제목
        if title:
            _add_text_box(slide, title,
                          x=1.5, y=3.9, w=10.33, h=1.1,
                          size=26, color=GRAY_LIGHT,
                          align=PP_ALIGN.CENTER)

        # 주차 범위
        if weeks:
            _add_text_box(slide, weeks,
                          x=1.5, y=5.0, w=10.33, h=0.6,
                          size=18, color=GRAY_MID,
                          align=PP_ALIGN.CENTER)

    def add_divider_slide(self, number: str, section_name: str) -> None:
        """섹션 구분 슬라이드 — 큰 번호 + 섹션명, 좌측 검정 세로 바."""
        slide = self._new_slide()

        # 좌측 검정 바
        bar = slide.shapes.add_shape(
            1,
            Inches(MARGIN), Inches(1.8),
            Inches(0.12), Inches(3.5),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = BLACK
        bar.line.fill.background()

        # 큰 주차 번호 (연회색)
        _add_text_box(slide, number,
                      x=MARGIN + 0.4, y=1.4, w=3.0, h=2.2,
                      size=100, bold=True, color=GRAY_LIGHT,
                      align=PP_ALIGN.LEFT)

        # 섹션명
        _add_text_box(slide, section_name,
                      x=MARGIN + 0.4, y=3.7, w=CONTENT_W - 0.4, h=1.4,
                      size=42, bold=True, color=BLACK)

    def save(self, filename: str = "") -> str:
        OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_name = "".join(c if c.isalnum() or c in " _-" else "_" for c in filename)
        fname = (f"{timestamp}_{safe_name[:40].strip()}.pptx"
                 if safe_name else f"{timestamp}_presentation.pptx")
        path = OUTPUT_DIR / fname
        self.prs.save(str(path))
        return str(path)
